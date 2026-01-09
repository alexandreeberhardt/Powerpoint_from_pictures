import os
from pptx import Presentation
from PIL import Image

def add_folder_to_presentation(prs, folder_path):
    valid_extensions = ('.jpg', '.jpeg', '.png', '.bmp')
    
    files = sorted([f for f in os.listdir(folder_path) if f.lower().endswith(valid_extensions)])

    slide_ratio = prs.slide_width / prs.slide_height

    for filename in files:
        file_path = os.path.join(folder_path, filename)
        
        # Layout 6 is an empty slide (Blank)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        with Image.open(file_path) as img:
            width_px, height_px = img.size

        img_ratio = width_px / height_px

        if img_ratio > slide_ratio:
            img_width = prs.slide_width
            img_height = img_width / img_ratio
        else:
            img_height = prs.slide_height
            img_width = img_height * img_ratio

        # centering
        left = (prs.slide_width - img_width) / 2
        top = (prs.slide_height - img_height) / 2

        slide.shapes.add_picture(file_path, left, top, width=img_width, height=img_height)

if __name__ == "__main__":
    output_filename = 'presentation.pptx'
    source_folders = [r".\dossier1", r".\dossier2", r".\dossier3"]
    
    prs = Presentation()

    for folder in source_folders:
        if os.path.exists(folder):
            add_folder_to_presentation(prs, folder)
        else:
            print(f"Warning: The folder {folder} is missing.")

    prs.save(output_filename)